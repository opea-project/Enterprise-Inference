"""
Document Parser
Handles extraction of text from various document formats
"""

import logging
from pathlib import Path
from typing import Dict, List, Optional
from pypdf import PdfReader  # Modern replacement for PyPDF2
from docx import Document
import openpyxl
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentParser:
    """
    Parse documents and extract text content.
    
    Supports multiple file formats (PDF, DOCX, XLSX, PPTX, TXT) and 
    extracts text into page/section dictionaries.
    """
    
    @staticmethod
    def parse_pdf(file_path: Path) -> Dict[int, str]:
        """
        Parse PDF and extract text by page.
        
        Args:
            file_path (Path): Path to the PDF file.
            
        Returns:
            Dict[int, str]: Dictionary mapping page numbers (1-based) to text content.
            
        Raises:
            ValueError: If parsing fails.
        """
        try:
            pages = {}
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        pages[page_num] = text
            
            logger.info(f"Extracted {len(pages)} pages from PDF: {file_path.name}")
            return pages
            
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            raise ValueError(f"Failed to parse PDF: {str(e)}")
    
    @staticmethod
    def parse_docx(file_path: Path) -> Dict[int, str]:
        """
        Parse DOCX and extract text by paragraph groupings (sections).
        
        Args:
            file_path (Path): Path to the DOCX file.
            
        Returns:
            Dict[int, str]: Dictionary mapping section numbers to text content.
            
        Raises:
            ValueError: If parsing fails.
        """
        try:
            doc = Document(file_path)
            sections = {}
            section_num = 1
            current_text = []
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    current_text.append(text)
                    # Group paragraphs into sections (every ~10 paragraphs)
                    if len(current_text) >= 10:
                        sections[section_num] = "\n".join(current_text)
                        current_text = []
                        section_num += 1
            
            # Add remaining text
            if current_text:
                sections[section_num] = "\n".join(current_text)
            
            logger.info(f"Extracted {len(sections)} sections from DOCX: {file_path.name}")
            return sections
            
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            raise ValueError(f"Failed to parse DOCX: {str(e)}")
    
    @staticmethod
    def parse_xlsx(file_path: Path) -> Dict[int, str]:
        """
        Parse XLSX and extract text from sheets.
        
        Args:
            file_path (Path): Path to the XLSX file.
            
        Returns:
            Dict[int, str]: Dictionary mapping sheet numbers (1-based) to text content.
            
        Raises:
            ValueError: If parsing fails.
        """
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            sheets = {}
            
            for sheet_num, sheet_name in enumerate(workbook.sheetnames, 1):
                sheet = workbook[sheet_name]
                rows = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        rows.append(row_text)
                
                if rows:
                    sheets[sheet_num] = f"Sheet: {sheet_name}\n" + "\n".join(rows)
            
            logger.info(f"Extracted {len(sheets)} sheets from XLSX: {file_path.name}")
            return sheets
            
        except Exception as e:
            logger.error(f"Error parsing XLSX {file_path}: {e}")
            raise ValueError(f"Failed to parse XLSX: {str(e)}")
    
    @staticmethod
    def parse_pptx(file_path: Path) -> Dict[int, str]:
        """
        Parse PPTX and extract text from slides.
        
        Args:
            file_path (Path): Path to the PPTX file.
            
        Returns:
            Dict[int, str]: Dictionary mapping slide numbers (1-based) to text content.
            
        Raises:
            ValueError: If parsing fails.
        """
        try:
            prs = Presentation(file_path)
            slides = {}
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)
                
                if text_parts:
                    slides[slide_num] = "\n".join(text_parts)
            
            logger.info(f"Extracted {len(slides)} slides from PPTX: {file_path.name}")
            return slides
            
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            raise ValueError(f"Failed to parse PPTX: {str(e)}")
    
    @staticmethod
    def parse_txt(file_path: Path) -> Dict[int, str]:
        """
        Parse TXT file.
        
        Args:
            file_path (Path): Path to the TXT file.
            
        Returns:
            Dict[int, str]: Dictionary with a single entry {1: text_content}.
            
        Raises:
            ValueError: If parsing fails.
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            if text.strip():
                return {1: text}
            return {}
            
        except Exception as e:
            logger.error(f"Error parsing TXT {file_path}: {e}")
            raise ValueError(f"Failed to parse TXT: {str(e)}")
    
    def parse_document(self, file_path: Path, file_type: str) -> Dict[int, str]:
        """
        Parse document based on file type.
        
        Dispatcher method that calls the appropriate parser.
        
        Args:
            file_path (Path): Path to document file.
            file_type (str): Type of file ('pdf', 'docx', 'xlsx', 'pptx', 'txt').
            
        Returns:
            Dict[int, str]: Dictionary mapping page/section numbers to text content.
            
        Raises:
            ValueError: If file type is unsupported.
        """
        file_type = file_type.lower().strip()
        
        parsers = {
            'pdf': self.parse_pdf,
            'docx': self.parse_docx,
            'xlsx': self.parse_xlsx,
            'pptx': self.parse_pptx,
            'ppt': self.parse_pptx,  # Alias
            'txt': self.parse_txt
        }
        
        parser = parsers.get(file_type)
        if not parser:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return parser(file_path)

